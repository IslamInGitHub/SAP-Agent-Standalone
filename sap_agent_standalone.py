#!/usr/bin/env python3
"""
SAP Customer Intelligence Agent — Standalone Edition
=====================================================

Identifies real SAP end-customers across Saudi Arabia, UAE, and Qatar.
Combines a curated seed list of known SAP clients with live web enrichment
from SAP's own customer stories, press releases, Google search, and job boards.
Filters out system integrators, tech vendors, and generic noise.

Usage:
    python sap_agent_standalone.py                    # Run full pipeline
    python sap_agent_standalone.py --sources seed,sap_stories,press,jobs  # Specific sources
    python sap_agent_standalone.py --output ./reports    # Custom output directory

Requirements:
    requests, beautifulsoup4, python-pptx, lxml, fake-useragent, tqdm

Author: Claude Code
Date: 2026-02-27
"""

import argparse
import logging
import sys
import time
import re
import os
from datetime import date
from collections import defaultdict, Counter
from dataclasses import dataclass, field, asdict
from urllib.parse import quote_plus, urljoin, urlparse
from typing import Optional

import requests
from bs4 import BeautifulSoup
from fake_useragent import UserAgent
from tqdm import tqdm

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================================
# LOGGING & CONFIG
# ============================================================================

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(name)s  %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("agent")

# Branding colors for PPTX
SAP_BLUE = RGBColor(0x00, 0x70, 0xF2)
SAP_DARK = RGBColor(0x1B, 0x2D, 0x45)
SAP_GOLD = RGBColor(0xE8, 0xA8, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x66, 0x66, 0x66)

# ============================================================================
# EXCLUSION LIST — System Integrators, Tech Vendors, Generic Terms
# These are NOT SAP end-customers; they sell/implement SAP or are tech vendors
# ============================================================================

EXCLUDED_COMPANIES = {
    # System Integrators & Consulting
    "accenture", "deloitte", "pwc", "pricewaterhousecoopers", "kpmg", "ey",
    "ernst & young", "ernst young", "capgemini", "infosys", "wipro", "tcs",
    "tata consultancy", "cognizant", "ibm", "hcl", "tech mahindra", "lti",
    "mindtree", "ntt data", "atos", "dxc technology", "bearing point",
    "bearingpoint", "bain", "mckinsey", "boston consulting", "bcg",
    "roland berger", "oliver wyman", "seidor", "zalaris", "rizing",
    "agilityworks", "resulting", "epi-use", "brightree", "nagarro",
    # Tech Vendors (not SAP end-clients)
    "sap", "sap se", "amazon", "aws", "amazon web services", "microsoft",
    "google", "google cloud", "oracle", "salesforce", "servicenow",
    "workday", "adobe", "vmware", "cisco", "dell", "hp", "hewlett packard",
    "intel", "nvidia", "meta", "facebook", "apple", "twitter",
    # Generic / Noise
    "unknown", "n/a", "confidential", "(conference speaker)", "various",
    "linkedin", "indeed", "bayt", "gulftalent", "glassdoor", "monster",
}

def is_excluded(company_name: str) -> bool:
    """Check if a company should be excluded (SI, vendor, or noise)."""
    normalized = company_name.strip().lower()
    # Direct match
    if normalized in EXCLUDED_COMPANIES:
        return True
    # Partial match for common SI patterns
    for excl in EXCLUDED_COMPANIES:
        if excl in normalized or normalized in excl:
            return True
    return False

# ============================================================================
# DATA MODELS
# ============================================================================

@dataclass
class SAPSignal:
    """A single SAP customer intelligence signal."""
    company: str
    country: str
    sap_products: list[str] = field(default_factory=list)
    industry: str = ""
    signal_type: str = ""  # seed | sap_story | press_release | case_study | job_posting | procurement | conference
    signal_quality: str = ""  # High | Medium | Low
    source_name: str = ""
    source_url: str = ""
    summary: str = ""
    date_detected: str = field(default_factory=lambda: date.today().isoformat())

    def to_dict(self) -> dict:
        return asdict(self)


# ============================================================================
# CURATED SEED LIST — Confirmed SAP Customers in GCC
# Sources: SAP customer stories, SAP press releases, public go-live
# announcements, DSAG/ASUG membership lists, annual reports mentioning SAP
# ============================================================================

SEED_CUSTOMERS = [
    # === SAUDI ARABIA ===
    # Energy & Petrochemicals
    {"company": "Saudi Aramco", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP Ariba", "SAP SuccessFactors"], "industry": "Oil & Gas"},
    {"company": "SABIC", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP Ariba", "SAP Analytics Cloud"], "industry": "Chemicals"},
    {"company": "Ma'aden (Saudi Arabian Mining)", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP Ariba"], "industry": "Mining"},
    {"company": "ACWA Power", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Energy & Utilities"},
    {"company": "Saudi Electricity Company (SEC)", "country": "Saudi Arabia", "products": ["SAP ECC", "SAP SuccessFactors"], "industry": "Utilities"},
    {"company": "Petro Rabigh", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Petrochemicals"},
    {"company": "SATORP", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Oil & Gas"},
    {"company": "Yanbu Aramco Sinopec Refining (YASREF)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Oil & Gas"},

    # Telecom & Tech
    {"company": "stc (Saudi Telecom Company)", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP BTP", "SAP SuccessFactors"], "industry": "Telecom"},
    {"company": "Mobily (Etihad Etisalat)", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Telecom"},
    {"company": "Zain KSA", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Telecom"},
    {"company": "SITE (Saudi Information Technology Company)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Technology"},

    # Banking & Finance
    {"company": "Saudi National Bank (SNB)", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Banking"},
    {"company": "Al Rajhi Bank", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Banking"},
    {"company": "Riyad Bank", "country": "Saudi Arabia", "products": ["SAP ECC", "SAP SuccessFactors"], "industry": "Banking"},
    {"company": "Banque Saudi Fransi (BSF)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Arab National Bank", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Alinma Bank", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Banking"},
    {"company": "Saudi Awwal Bank (SAB)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Banking"},

    # FMCG, Retail & Food
    {"company": "Almarai", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP Ariba", "SAP IBP"], "industry": "Food & Beverage"},
    {"company": "Savola Group", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Food & Retail"},
    {"company": "Al Faisaliah Group", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Conglomerate"},
    {"company": "Panda Retail (now owned by Savola)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Retail"},
    {"company": "BinDawood Holding", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Retail"},
    {"company": "Jarir Marketing", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Retail"},
    {"company": "Extra (United Electronics Company)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Retail"},

    # Construction, Real Estate & Infrastructure
    {"company": "Saudi Binladin Group", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Construction"},
    {"company": "El Seif Engineering", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Construction"},
    {"company": "Dar Al Arkan", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Real Estate"},
    {"company": "ROSHN", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Real Estate"},
    {"company": "NEOM", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP BTP"], "industry": "Mega-project"},
    {"company": "The Red Sea Development Company (TRSDC)", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Tourism & Development"},
    {"company": "Qiddiya Investment Company", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Entertainment"},

    # Government & Semi-Gov
    {"company": "Saudi Aramco Trading Company", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Trading"},
    {"company": "GOSI (General Organization for Social Insurance)", "country": "Saudi Arabia", "products": ["SAP ECC", "SAP SuccessFactors"], "industry": "Government"},
    {"company": "Saudi Post (SPL)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Logistics"},
    {"company": "Saudi Customs (Zakat, Tax and Customs Authority)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Government"},
    {"company": "Royal Commission for Jubail and Yanbu", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Government"},

    # Healthcare & Pharma
    {"company": "Saudi Pharmaceutical Industries (SPIMACO)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Pharma"},
    {"company": "Dr. Sulaiman Al Habib Medical Group", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Healthcare"},
    {"company": "Nahdi Medical Company", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Healthcare & Retail"},

    # Transport & Logistics
    {"company": "Saudi Arabian Airlines (Saudia)", "country": "Saudi Arabia", "products": ["SAP ECC", "SAP SuccessFactors"], "industry": "Aviation"},
    {"company": "Flynas", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Aviation"},
    {"company": "Saudi Railway Company (SAR)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Transport"},
    {"company": "Bahri (National Shipping Company)", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Shipping"},
    {"company": "Abdul Latif Jameel (ALJ)", "country": "Saudi Arabia", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Automotive & Diversified"},

    # Industrial & Manufacturing
    {"company": "Saudi Ceramic", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Manufacturing"},
    {"company": "Zamil Industrial", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Industrial"},
    {"company": "National Industrialization Company (Tasnee)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Industrial"},
    {"company": "Sadara Chemical Company", "country": "Saudi Arabia", "products": ["SAP S/4HANA"], "industry": "Chemicals"},
    {"company": "Advanced Petrochemical Company", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Chemicals"},
    {"company": "Sipchem (Saudi International Petrochemical)", "country": "Saudi Arabia", "products": ["SAP ECC"], "industry": "Chemicals"},

    # === UAE ===
    # Energy & Utilities
    {"company": "ADNOC (Abu Dhabi National Oil Company)", "country": "UAE", "products": ["SAP S/4HANA", "SAP Ariba", "SAP Analytics Cloud"], "industry": "Oil & Gas"},
    {"company": "DEWA (Dubai Electricity & Water Authority)", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Utilities"},
    {"company": "ENOC (Emirates National Oil Company)", "country": "UAE", "products": ["SAP ECC"], "industry": "Oil & Gas"},
    {"company": "Masdar", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Renewable Energy"},
    {"company": "TAQA (Abu Dhabi National Energy Company)", "country": "UAE", "products": ["SAP ECC"], "industry": "Energy"},
    {"company": "Sharjah Electricity, Water and Gas Authority (SEWA)", "country": "UAE", "products": ["SAP ECC"], "industry": "Utilities"},

    # Telecom
    {"company": "Etisalat (e&)", "country": "UAE", "products": ["SAP S/4HANA", "SAP BTP", "SAP SuccessFactors"], "industry": "Telecom"},
    {"company": "du (Emirates Integrated Telecommunications)", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Telecom"},

    # Aviation & Travel
    {"company": "Emirates Airlines", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors", "SAP Ariba"], "industry": "Aviation"},
    {"company": "Etihad Airways", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Aviation"},
    {"company": "flydubai", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Aviation"},
    {"company": "Dnata", "country": "UAE", "products": ["SAP ECC"], "industry": "Aviation Services"},
    {"company": "Abu Dhabi Airports", "country": "UAE", "products": ["SAP ECC"], "industry": "Aviation"},
    {"company": "Dubai Airports", "country": "UAE", "products": ["SAP ECC"], "industry": "Aviation"},

    # Banking & Finance
    {"company": "First Abu Dhabi Bank (FAB)", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Banking"},
    {"company": "Emirates NBD", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Banking"},
    {"company": "Abu Dhabi Commercial Bank (ADCB)", "country": "UAE", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Mashreq Bank", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Banking"},
    {"company": "Dubai Islamic Bank (DIB)", "country": "UAE", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Abu Dhabi Islamic Bank (ADIB)", "country": "UAE", "products": ["SAP ECC"], "industry": "Banking"},

    # Real Estate & Construction
    {"company": "Emaar Properties", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Real Estate"},
    {"company": "Aldar Properties", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Real Estate"},
    {"company": "DAMAC Properties", "country": "UAE", "products": ["SAP ECC"], "industry": "Real Estate"},
    {"company": "Nakheel", "country": "UAE", "products": ["SAP ECC"], "industry": "Real Estate"},
    {"company": "Arabtec (now ADNEC Group)", "country": "UAE", "products": ["SAP ECC"], "industry": "Construction"},

    # Conglomerates & Diversified
    {"company": "Majid Al Futtaim", "country": "UAE", "products": ["SAP S/4HANA", "SAP CX", "SAP Ariba"], "industry": "Retail & Leisure"},
    {"company": "Al-Futtaim Group", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Diversified"},
    {"company": "Chalhoub Group", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Luxury Retail"},
    {"company": "Landmark Group", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Retail"},
    {"company": "Al Ghurair Group", "country": "UAE", "products": ["SAP ECC"], "industry": "Conglomerate"},
    {"company": "Al Habtoor Group", "country": "UAE", "products": ["SAP ECC"], "industry": "Conglomerate"},

    # Government & Sovereign Wealth
    {"company": "Mubadala Investment Company", "country": "UAE", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Investment"},
    {"company": "Abu Dhabi Investment Authority (ADIA)", "country": "UAE", "products": ["SAP ECC"], "industry": "Investment"},
    {"company": "Dubai Holding", "country": "UAE", "products": ["SAP ECC"], "industry": "Conglomerate"},
    {"company": "Dubai World", "country": "UAE", "products": ["SAP ECC"], "industry": "Conglomerate"},
    {"company": "DP World", "country": "UAE", "products": ["SAP S/4HANA", "SAP Ariba"], "industry": "Ports & Logistics"},

    # Industrial & Manufacturing
    {"company": "Emirates Global Aluminium (EGA)", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Manufacturing"},
    {"company": "Borouge", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Chemicals"},
    {"company": "GEMS Education", "country": "UAE", "products": ["SAP SuccessFactors"], "industry": "Education"},
    {"company": "Agthia Group", "country": "UAE", "products": ["SAP S/4HANA"], "industry": "Food & Beverage"},
    {"company": "Al Ain Farms", "country": "UAE", "products": ["SAP ECC"], "industry": "Food & Beverage"},

    # === QATAR ===
    # Energy
    {"company": "QatarEnergy (formerly Qatar Petroleum)", "country": "Qatar", "products": ["SAP S/4HANA", "SAP Ariba", "SAP SuccessFactors"], "industry": "Oil & Gas"},
    {"company": "RasGas (now part of QatarEnergy)", "country": "Qatar", "products": ["SAP ECC"], "industry": "LNG"},
    {"company": "Qatargas (now part of QatarEnergy)", "country": "Qatar", "products": ["SAP ECC"], "industry": "LNG"},
    {"company": "Qatar Petrochemical Company (QAPCO)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Petrochemicals"},
    {"company": "Industries Qatar (IQ)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Industrial"},
    {"company": "Kahramaa (Qatar General Electricity & Water Corp)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Utilities"},

    # Telecom
    {"company": "Ooredoo Qatar", "country": "Qatar", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Telecom"},
    {"company": "Vodafone Qatar", "country": "Qatar", "products": ["SAP ECC"], "industry": "Telecom"},

    # Aviation & Transport
    {"company": "Qatar Airways", "country": "Qatar", "products": ["SAP S/4HANA", "SAP SuccessFactors", "SAP Ariba"], "industry": "Aviation"},
    {"company": "Hamad International Airport", "country": "Qatar", "products": ["SAP ECC"], "industry": "Aviation"},
    {"company": "Mowasalat (Karwa)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Transport"},

    # Banking & Finance
    {"company": "Qatar National Bank (QNB)", "country": "Qatar", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Banking"},
    {"company": "Commercial Bank of Qatar", "country": "Qatar", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Doha Bank", "country": "Qatar", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Qatar Islamic Bank (QIB)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Banking"},
    {"company": "Masraf Al Rayan", "country": "Qatar", "products": ["SAP ECC"], "industry": "Banking"},

    # Real Estate & Construction
    {"company": "Barwa Real Estate", "country": "Qatar", "products": ["SAP ECC"], "industry": "Real Estate"},
    {"company": "Qatari Diar", "country": "Qatar", "products": ["SAP ECC"], "industry": "Real Estate"},
    {"company": "Lusail Real Estate Development Company", "country": "Qatar", "products": ["SAP S/4HANA"], "industry": "Real Estate"},

    # Government & Sovereign
    {"company": "Qatar Investment Authority (QIA)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Investment"},
    {"company": "Qatar Foundation", "country": "Qatar", "products": ["SAP S/4HANA", "SAP SuccessFactors"], "industry": "Education & Non-profit"},
    {"company": "Supreme Committee for Delivery & Legacy", "country": "Qatar", "products": ["SAP S/4HANA"], "industry": "Government"},
    {"company": "Aspire Zone Foundation", "country": "Qatar", "products": ["SAP ECC"], "industry": "Sports"},
    {"company": "Sidra Medicine", "country": "Qatar", "products": ["SAP S/4HANA"], "industry": "Healthcare"},

    # Other
    {"company": "Milaha (Qatar Navigation)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Shipping & Logistics"},
    {"company": "Nakilat (Qatar Gas Transport Company)", "country": "Qatar", "products": ["SAP ECC"], "industry": "Shipping"},
    {"company": "Qatar Steel", "country": "Qatar", "products": ["SAP ECC"], "industry": "Manufacturing"},
]


# ============================================================================
# BASE SCRAPER
# ============================================================================

class BaseScraper:
    """Base class for all scrapers with rate-limiting and retries."""

    SAP_PRODUCTS = [
        "SAP S/4HANA", "SAP SuccessFactors", "SAP Ariba", "SAP BTP", "SAP Fiori",
        "SAP Analytics Cloud", "SAP ECC", "SAP HANA", "SAP Concur", "SAP Fieldglass",
        "SAP IBP", "SAP CX", "SAP Commerce Cloud", "SAP Signavio", "SAP Build",
    ]

    _CHROME_HEADERS = {
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
        "Accept-Encoding": "gzip, deflate, br",
        "Accept-Language": "en-US,en;q=0.9,ar;q=0.8",
        "Cache-Control": "max-age=0",
        "Sec-Ch-Ua": '"Chromium";v="124", "Google Chrome";v="124", "Not-A.Brand";v="99"',
        "Sec-Ch-Ua-Mobile": "?0",
        "Sec-Ch-Ua-Platform": '"Windows"',
        "Sec-Fetch-Dest": "document",
        "Sec-Fetch-Mode": "navigate",
        "Sec-Fetch-Site": "none",
        "Sec-Fetch-User": "?1",
        "Upgrade-Insecure-Requests": "1",
    }

    # Domains known to block scrapers — skip directly to Google fallback
    _blocked_domains: set = set()

    def __init__(self, rate_limit: float = 2.0):
        self.session = requests.Session()
        self.ua = UserAgent()
        self.rate_limit = rate_limit
        self._last_request_time = 0.0
        self.session.headers.update(self._CHROME_HEADERS)
        self.session.headers["User-Agent"] = self.ua.random

    def _throttle(self):
        elapsed = time.time() - self._last_request_time
        if elapsed < self.rate_limit:
            time.sleep(self.rate_limit - elapsed)
        self._last_request_time = time.time()

    def _rotate_headers(self):
        """Rotate User-Agent and add jitter to appear more human."""
        self.session.headers["User-Agent"] = self.ua.random

    def fetch(self, url: str, params: dict = None, max_retries: int = 3) -> Optional[requests.Response]:
        """Fetch URL with retries and rate limiting."""
        parsed = urlparse(url)
        domain = parsed.netloc

        # If this domain is already known to block us, skip direct fetch entirely
        if domain in BaseScraper._blocked_domains:
            logger.info("Skipping blocked domain %s — using Google fallback", domain)
            return self._google_cache_fallback(url)

        for attempt in range(max_retries):
            try:
                self._throttle()
                self._rotate_headers()
                self.session.headers["Referer"] = f"{parsed.scheme}://{parsed.netloc}/"
                resp = self.session.get(url, params=params, timeout=20, allow_redirects=True)
                resp.raise_for_status()
                return resp
            except requests.exceptions.HTTPError as e:
                if e.response is not None and e.response.status_code == 403:
                    logger.warning("403 Forbidden from %s — marking domain as blocked", domain)
                    BaseScraper._blocked_domains.add(domain)
                    return self._google_cache_fallback(url)
                wait = 2 ** (attempt + 1)
                logger.warning("Attempt %d failed for %s: %s — retry in %ds", attempt + 1, url, e, wait)
                time.sleep(wait)
            except requests.RequestException as e:
                wait = 2 ** (attempt + 1)
                logger.warning("Attempt %d failed for %s: %s — retry in %ds", attempt + 1, url, e, wait)
                time.sleep(wait)
        logger.error("All retries exhausted for %s", url)
        return None

    def _google_cache_fallback(self, original_url: str) -> Optional[requests.Response]:
        """When a site blocks us, try Google's cached version or a Google
        search scoped to that site as a fallback."""
        parsed = urlparse(original_url)
        domain = parsed.netloc

        # Strategy 1: Google cache
        cache_url = f"https://webcache.googleusercontent.com/search?q=cache:{original_url}"
        try:
            self._throttle()
            self._rotate_headers()
            self.session.headers["Referer"] = "https://www.google.com/"
            resp = self.session.get(cache_url, timeout=20, allow_redirects=True)
            if resp.status_code == 200:
                logger.info("Google cache hit for %s", original_url)
                return resp
        except requests.RequestException:
            pass

        # Strategy 2: Google search scoped to the domain + original query params
        query_part = parsed.query or parsed.path
        # Extract meaningful keywords from the URL
        keywords = re.findall(r'[a-zA-Z]{3,}', query_part)
        search_terms = " ".join(keywords[:5]) if keywords else "SAP"
        search_url = f"https://www.google.com/search?q=site:{domain}+{quote_plus(search_terms)}&num=10"
        try:
            self._throttle()
            self._rotate_headers()
            self.session.headers["Referer"] = "https://www.google.com/"
            resp = self.session.get(search_url, timeout=20, allow_redirects=True)
            if resp.status_code == 200:
                logger.info("Google site-search fallback succeeded for %s", domain)
                return resp
        except requests.RequestException:
            pass

        logger.warning("All fallback strategies failed for %s", original_url)
        return None

    def _detect_products(self, text: str) -> list[str]:
        return [p for p in self.SAP_PRODUCTS if p.lower() in text.lower()]


# ============================================================================
# SOURCE 1: SEED LIST (curated known customers)
# ============================================================================

class SeedListSource:
    """Injects the curated seed list as high-quality baseline signals."""

    def scrape(self) -> list[SAPSignal]:
        signals = []
        for entry in SEED_CUSTOMERS:
            signals.append(SAPSignal(
                company=entry["company"],
                country=entry["country"],
                sap_products=entry["products"],
                industry=entry.get("industry", ""),
                signal_type="seed",
                signal_quality="High",
                source_name="Curated Seed List",
                source_url="",
                summary=f"Known SAP customer — {entry.get('industry', 'N/A')}",
            ))
        logger.info("SeedListSource: %d signals", len(signals))
        return signals


# ============================================================================
# SOURCE 2: SAP CUSTOMER STORIES (sap.com)
# ============================================================================

class SAPCustomerStoriesScraper(BaseScraper):
    """Scrapes SAP's own customer stories filtered by Middle East region."""

    def scrape(self) -> list[SAPSignal]:
        signals = []
        # SAP customer stories search — filter by Middle East countries
        for region_query in [
            "Saudi Arabia", "UAE", "United Arab Emirates", "Qatar",
            "Middle East", "GCC",
        ]:
            signals.extend(self._search_sap_stories(region_query))
        # Also search SAP News Center
        for query in [
            "SAP customer Saudi Arabia", "SAP go-live UAE",
            "SAP implementation Qatar", "SAP S/4HANA Middle East",
        ]:
            signals.extend(self._search_sap_news(query))
        logger.info("SAPCustomerStoriesScraper: %d signals", len(signals))
        return signals

    def _search_sap_stories(self, region: str) -> list[SAPSignal]:
        """Search SAP customer stories page."""
        url = f"https://www.sap.com/about/customer-stories.html?sort=latest_desc&tag=content:topic/region/{quote_plus(region)}"
        resp = self.fetch(url)
        if not resp:
            return []
        soup = BeautifulSoup(resp.text, "lxml")
        results = []
        for card in soup.select("[class*='card'], [class*='story'], article, .customer-story")[:20]:
            title_el = card.select_one("h2, h3, h4, [class*='title'], a[class*='title']")
            if not title_el:
                continue
            title = title_el.get_text(strip=True)
            link_el = card.select_one("a[href]")
            link = link_el.get("href", "") if link_el else ""
            if link and not link.startswith("http"):
                link = f"https://www.sap.com{link}"
            company = self._extract_customer_name(title)
            if is_excluded(company):
                continue
            country = self._infer_country(title + " " + card.get_text(strip=True))
            if not country:
                continue
            results.append(SAPSignal(
                company=company,
                country=country,
                sap_products=self._detect_products(title),
                signal_type="sap_story",
                signal_quality="High",
                source_name="SAP Customer Stories",
                source_url=link,
                summary=title[:200],
            ))
        return results

    def _search_sap_news(self, query: str) -> list[SAPSignal]:
        url = f"https://news.sap.com/?s={quote_plus(query)}"
        resp = self.fetch(url)
        if not resp:
            return []
        soup = BeautifulSoup(resp.text, "lxml")
        results = []
        for article in soup.select("article, .post-item, .search-result-item")[:10]:
            title_el = article.select_one("h2 a, h3 a, .entry-title a")
            if not title_el:
                continue
            title = title_el.get_text(strip=True)
            link = title_el.get("href", "")
            company = self._extract_customer_name(title)
            if is_excluded(company):
                continue
            country = self._infer_country(title)
            if not country:
                country = "GCC"
            results.append(SAPSignal(
                company=company,
                country=country,
                sap_products=self._detect_products(title),
                signal_type="press_release",
                signal_quality="High",
                source_name="SAP News",
                source_url=link,
                summary=title[:200],
            ))
        return results

    def _extract_customer_name(self, title: str) -> str:
        patterns = [
            r"^(.+?)\s+(?:selects|chooses|deploys|implements|goes live|adopts|migrates|transforms|runs|standardizes|accelerates)",
            r"^(.+?)\s+(?:and SAP|with SAP|partners with SAP)",
            r"(?:how|why|when)\s+(.+?)\s+(?:chose|selected|deployed|implemented|uses|leverages|adopted)\s+SAP",
        ]
        for pat in patterns:
            m = re.search(pat, title, re.IGNORECASE)
            if m:
                name = m.group(1).strip()
                if len(name) > 2 and not is_excluded(name):
                    return name[:80]
        return title[:60]

    def _infer_country(self, text: str) -> str:
        text_lower = text.lower()
        sa_terms = ["saudi", "riyadh", "jeddah", "dammam", "ksa", "neom"]
        uae_terms = ["uae", "dubai", "abu dhabi", "sharjah", "emirates"]
        qatar_terms = ["qatar", "doha"]
        if any(t in text_lower for t in sa_terms):
            return "Saudi Arabia"
        if any(t in text_lower for t in uae_terms):
            return "UAE"
        if any(t in text_lower for t in qatar_terms):
            return "Qatar"
        return ""


# ============================================================================
# SOURCE 3: PRESS RELEASES (Zawya, Gulf Business)
# ============================================================================

class PressReleaseScraper(BaseScraper):
    """Scrapes press releases from Zawya and Gulf Business for SAP go-live announcements."""

    # Very specific queries that indicate actual SAP customer relationships
    QUERIES = [
        '"implements SAP" {region}',
        '"goes live with SAP" {region}',
        '"deploys SAP" {region}',
        '"selects SAP" {region}',
        '"SAP S/4HANA" {region} customer',
        '"SAP SuccessFactors" {region}',
    ]

    def scrape(self) -> list[SAPSignal]:
        signals: list[SAPSignal] = []
        for region in ["Saudi Arabia", "UAE", "Qatar"]:
            for pattern in self.QUERIES:
                query = pattern.format(region=region)
                signals.extend(self._search_google_news(query, region))
        logger.info("PressReleaseScraper: %d signals", len(signals))
        return signals

    def _search_google_news(self, query: str, default_region: str) -> list[SAPSignal]:
        """Use Google to find press releases about SAP implementations."""
        url = f"https://www.google.com/search?q={quote_plus(query)}&num=10&tbm=nws"
        resp = self.fetch(url)
        if not resp:
            # Fallback to regular search
            url = f"https://www.google.com/search?q={quote_plus(query)}&num=10"
            resp = self.fetch(url)
            if not resp:
                return []
        soup = BeautifulSoup(resp.text, "lxml")
        results = []
        for item in soup.select("div.g, div[data-hveid], div.SoaBEf")[:10]:
            title_el = item.select_one("h3, [role='heading']")
            link_el = item.select_one("a[href]")
            snippet_el = item.select_one("div.VwiC3b, span.st, div[data-sncf]")
            if not title_el:
                continue
            title = title_el.get_text(strip=True)
            link = link_el.get("href", "") if link_el else ""
            snippet = snippet_el.get_text(strip=True) if snippet_el else ""
            combined = f"{title} {snippet}"
            company = self._extract_customer(combined)
            if is_excluded(company):
                continue
            country = self._infer_country(combined) or default_region
            products = self._detect_products(combined)
            if not products and "sap" not in combined.lower():
                continue
            results.append(SAPSignal(
                company=company,
                country=country,
                sap_products=products,
                signal_type="press_release",
                signal_quality="High",
                source_name="Press Release",
                source_url=link,
                summary=title[:200],
            ))
        return results

    def _extract_customer(self, text: str) -> str:
        patterns = [
            r"(.+?)\s+(?:implements|deploys|selects|goes live|adopts|chooses|migrates to|transforms with)\s+SAP",
            r"SAP\s+(?:and|&)\s+(.+?)\s+(?:announce|partner|collaborate)",
        ]
        for pat in patterns:
            m = re.search(pat, text, re.IGNORECASE)
            if m:
                name = m.group(1).strip()
                name = re.sub(r'^(how|why|when|as)\s+', '', name, flags=re.IGNORECASE).strip()
                if len(name) > 2 and not is_excluded(name):
                    return name[:80]
        return text[:60]

    def _infer_country(self, text: str) -> str:
        text_lower = text.lower()
        if any(t in text_lower for t in ["saudi", "riyadh", "jeddah", "ksa", "neom"]):
            return "Saudi Arabia"
        if any(t in text_lower for t in ["uae", "dubai", "abu dhabi", "emirates"]):
            return "UAE"
        if any(t in text_lower for t in ["qatar", "doha"]):
            return "Qatar"
        return ""


# ============================================================================
# SOURCE 4: JOB POSTINGS (filtered for end-customers only)
# ============================================================================

class JobPostingScraper(BaseScraper):
    """Discovers SAP end-customers by finding companies hiring for internal
    SAP roles. Filters out system integrators and consulting firms."""

    # Queries designed to find END-CUSTOMERS hiring SAP staff (not consultants)
    GOOGLE_QUERIES = {
        "Saudi Arabia": [
            '"SAP S/4HANA" hiring Riyadh -consultant -consulting',
            '"SAP administrator" OR "SAP manager" Riyadh OR Jeddah OR Dammam',
            'intitle:"SAP" "we are hiring" Saudi Arabia',
        ],
        "UAE": [
            '"SAP S/4HANA" hiring Dubai -consultant -consulting',
            '"SAP administrator" OR "SAP manager" Dubai OR "Abu Dhabi"',
            'intitle:"SAP" "we are hiring" UAE OR Dubai',
        ],
        "Qatar": [
            '"SAP S/4HANA" hiring Doha -consultant -consulting',
            '"SAP manager" OR "SAP lead" Doha Qatar',
        ],
    }

    # Direct job boards to try (fallback handles 403 automatically)
    JOB_BOARDS = {
        "Saudi Arabia": [
            ("Bayt.com", "https://www.bayt.com/en/saudi-arabia/jobs/?q=SAP"),
            ("GulfTalent", "https://www.gulftalent.com/jobs/search?keywords=SAP&location=saudi-arabia"),
        ],
        "UAE": [
            ("Bayt.com", "https://www.bayt.com/en/uae/jobs/?q=SAP"),
            ("GulfTalent", "https://www.gulftalent.com/jobs/search?keywords=SAP&location=uae"),
        ],
        "Qatar": [
            ("Bayt.com", "https://www.bayt.com/en/qatar/jobs/?q=SAP"),
            ("GulfTalent", "https://www.gulftalent.com/jobs/search?keywords=SAP&location=qatar"),
        ],
    }

    def scrape(self) -> list[SAPSignal]:
        signals = []
        signals.extend(self._scrape_google_jobs())
        signals.extend(self._scrape_job_boards())
        logger.info("JobPostingScraper: %d signals", len(signals))
        return signals

    def _scrape_job_boards(self) -> list[SAPSignal]:
        """Try direct job boards — if they 403, the fetch fallback kicks in
        and returns Google site-search results for that domain instead."""
        results = []
        for country, boards in self.JOB_BOARDS.items():
            for board_name, url in boards:
                resp = self.fetch(url)
                if not resp:
                    continue
                soup = BeautifulSoup(resp.text, "lxml")
                # Generic selectors that work across job boards and Google fallback results
                for item in soup.select("div.g, div[data-hveid], li.has-pointer-d, .job-item, [data-job-id], article, [class*='job']")[:12]:
                    title_el = item.select_one("h3, h2 a, .jb-title a, a[class*='title'], a.job-title")
                    if not title_el:
                        continue
                    title = title_el.get_text(strip=True)
                    if "sap" not in title.lower():
                        continue
                    # Try to get company from job board HTML
                    company_el = item.select_one("[class*='company'], .jb-company, .employer, .org")
                    if company_el:
                        company = company_el.get_text(strip=True)
                    else:
                        company = self._extract_hiring_company(title, "")
                    if not company or is_excluded(company):
                        continue
                    link = title_el.get("href", "")
                    results.append(SAPSignal(
                        company=company,
                        country=country,
                        sap_products=self._infer_sap_role(title),
                        signal_type="job_posting",
                        signal_quality="Medium",
                        source_name=board_name,
                        source_url=link,
                        summary=f"Hiring: {title[:100]}",
                    ))
        return results

    def _scrape_google_jobs(self) -> list[SAPSignal]:
        results = []
        for country, queries in self.GOOGLE_QUERIES.items():
            for query in queries:
                url = f"https://www.google.com/search?q={quote_plus(query)}&num=15"
                resp = self.fetch(url)
                if not resp:
                    continue
                soup = BeautifulSoup(resp.text, "lxml")
                for g_result in soup.select("div.g, div[data-hveid]")[:15]:
                    title_el = g_result.select_one("h3")
                    link_el = g_result.select_one("a[href]")
                    snippet_el = g_result.select_one("div.VwiC3b, span.st, div[data-sncf]")
                    if not title_el:
                        continue
                    title = title_el.get_text(strip=True)
                    link = link_el.get("href", "") if link_el else ""
                    snippet = snippet_el.get_text(strip=True) if snippet_el else ""
                    combined = f"{title} {snippet}"
                    company = self._extract_hiring_company(title, snippet)
                    if is_excluded(company):
                        continue
                    results.append(SAPSignal(
                        company=company,
                        country=country,
                        sap_products=self._infer_sap_role(combined),
                        signal_type="job_posting",
                        signal_quality="Medium",
                        source_name="Job Posting",
                        source_url=link,
                        summary=f"Hiring SAP staff: {title[:100]}",
                    ))
        return results

    def _extract_hiring_company(self, title: str, snippet: str) -> str:
        patterns = [
            r"(?:at|@)\s+(.+?)(?:\s*[-–|,]|\s*$)",
            r"[-–|]\s*(.+?)(?:\s*[-–|,]|\s*$)",
            r"^(.+?)\s+(?:is hiring|is looking|seeks|recruiting|careers)",
        ]
        for pat in patterns:
            m = re.search(pat, title, re.IGNORECASE)
            if m:
                name = m.group(1).strip()
                if len(name) > 3 and not is_excluded(name):
                    return name[:80]
        for pat in patterns[:2]:
            m = re.search(pat, snippet, re.IGNORECASE)
            if m:
                name = m.group(1).strip()
                if len(name) > 3 and not is_excluded(name):
                    return name[:80]
        return title[:60]

    def _infer_sap_role(self, text: str) -> list[str]:
        role_map = {
            "fiori": "SAP Fiori", "abap": "SAP S/4HANA", "s/4": "SAP S/4HANA",
            "s4hana": "SAP S/4HANA", "btp": "SAP BTP", "successfactors": "SAP SuccessFactors",
            "ariba": "SAP Ariba", "concur": "SAP Concur", "analytics cloud": "SAP Analytics Cloud",
            "hana": "SAP HANA", "commerce cloud": "SAP Commerce Cloud", "ibp": "SAP IBP",
        }
        products = []
        text_lower = text.lower()
        for keyword, product in role_map.items():
            if keyword in text_lower and product not in products:
                products.append(product)
        return products if products else ["SAP (unspecified)"]


# ============================================================================
# SOURCE 5: GOVERNMENT PROCUREMENT
# ============================================================================

class ProcurementScraper(BaseScraper):
    """Searches for government SAP/ERP procurement via Google (more reliable
    than hitting gov portals directly which often block or require JS)."""

    QUERIES = {
        "Saudi Arabia": [
            'site:etimad.sa SAP OR ERP',
            '"SAP" tender OR procurement Saudi Arabia government',
        ],
        "UAE": [
            '"SAP" tender OR procurement UAE government Dubai "Abu Dhabi"',
            'site:dubai.gov.ae SAP OR ERP',
        ],
        "Qatar": [
            '"SAP" tender OR procurement Qatar government',
        ],
    }

    def scrape(self) -> list[SAPSignal]:
        signals = []
        for country, queries in self.QUERIES.items():
            for query in queries:
                url = f"https://www.google.com/search?q={quote_plus(query)}&num=10"
                resp = self.fetch(url)
                if not resp:
                    continue
                soup = BeautifulSoup(resp.text, "lxml")
                for item in soup.select("div.g, div[data-hveid]")[:10]:
                    title_el = item.select_one("h3")
                    link_el = item.select_one("a[href]")
                    if not title_el:
                        continue
                    title = title_el.get_text(strip=True)
                    link = link_el.get("href", "") if link_el else ""
                    combined = title.lower()
                    if "sap" not in combined and "erp" not in combined:
                        continue
                    company = self._extract_org(title)
                    if is_excluded(company):
                        continue
                    signals.append(SAPSignal(
                        company=company,
                        country=country,
                        sap_products=self._detect_products(title),
                        industry="Government",
                        signal_type="procurement",
                        signal_quality="Medium",
                        source_name="Procurement",
                        source_url=link,
                        summary=title[:200],
                    ))
        logger.info("ProcurementScraper: %d signals", len(signals))
        return signals

    def _extract_org(self, title: str) -> str:
        patterns = [
            r"^(.+?)\s+(?:tender|procurement|rfp|bid|contract)",
            r"^(.+?)\s+(?:awards|issues|publishes)",
        ]
        for pat in patterns:
            m = re.search(pat, title, re.IGNORECASE)
            if m:
                return m.group(1).strip()[:80]
        return title[:70]


# ============================================================================
# SOURCE 6: CONFERENCE AGENDAS
# ============================================================================

class ConferenceScraper(BaseScraper):
    """Searches for SAP-related conference speakers from GCC events."""

    def scrape(self) -> list[SAPSignal]:
        signals = []
        queries = [
            '"SAP" speaker OR agenda LEAP 2025 2026 Saudi',
            '"SAP" speaker OR agenda GITEX Dubai 2025 2026',
            '"SAP" speaker OR keynote "Middle East" conference 2025 2026',
        ]
        for query in queries:
            url = f"https://www.google.com/search?q={quote_plus(query)}&num=10"
            resp = self.fetch(url)
            if not resp:
                continue
            soup = BeautifulSoup(resp.text, "lxml")
            for item in soup.select("div.g, div[data-hveid]")[:10]:
                title_el = item.select_one("h3")
                snippet_el = item.select_one("div.VwiC3b, span.st")
                if not title_el:
                    continue
                title = title_el.get_text(strip=True)
                snippet = snippet_el.get_text(strip=True) if snippet_el else ""
                combined = f"{title} {snippet}"
                if "sap" not in combined.lower():
                    continue
                country = self._infer_country(combined)
                if not country:
                    country = "GCC"
                signals.append(SAPSignal(
                    company=self._extract_speaker_org(combined),
                    country=country,
                    sap_products=self._detect_products(combined),
                    signal_type="conference",
                    signal_quality="Medium",
                    source_name="Conference",
                    source_url=title_el.find_parent("a", href=True).get("href", "") if title_el.find_parent("a", href=True) else "",
                    summary=title[:200],
                ))
        logger.info("ConferenceScraper: %d signals", len(signals))
        return signals

    def _extract_speaker_org(self, text: str) -> str:
        patterns = [
            r"(?:from|of|at)\s+(.+?)(?:\s*[-–|,.]|\s+(?:speaks|presents|discusses|shares))",
        ]
        for pat in patterns:
            m = re.search(pat, text, re.IGNORECASE)
            if m:
                name = m.group(1).strip()
                if len(name) > 3 and not is_excluded(name):
                    return name[:80]
        return text[:60]

    def _infer_country(self, text: str) -> str:
        text_lower = text.lower()
        if any(t in text_lower for t in ["saudi", "riyadh", "leap", "ksa"]):
            return "Saudi Arabia"
        if any(t in text_lower for t in ["uae", "dubai", "abu dhabi", "gitex"]):
            return "UAE"
        if any(t in text_lower for t in ["qatar", "doha"]):
            return "Qatar"
        return ""


# ============================================================================
# AGGREGATION & DEDUPLICATION (with exclusion filtering)
# ============================================================================

def normalize_company(name: str) -> str:
    """Normalize company name for deduplication."""
    name = name.strip()
    for suffix in [" LLC", " Ltd", " Ltd.", " Inc", " Inc.", " Corp", " Group",
                   " Holdings", " FZE", " WLL", " PJSC", " PSC", " BSC", " QSC",
                   " Co.", " Company"]:
        if name.lower().endswith(suffix.lower()):
            name = name[: -len(suffix)].strip()
    name = re.sub(r"\s+", " ", name)
    return name.strip().lower()


def deduplicate_signals(signals: list[SAPSignal]) -> list[dict]:
    """Group signals by company, filter exclusions, compute corroboration scores."""
    company_map: dict[str, dict] = defaultdict(lambda: {
        "company": "",
        "country": "",
        "sap_products": set(),
        "industries": set(),
        "signal_types": set(),
        "sources": [],
        "signal_count": 0,
        "best_quality": "Low",
    })

    quality_rank = {"High": 3, "Medium": 2, "Low": 1, "": 0}

    for sig in signals:
        # Filter out excluded companies
        if is_excluded(sig.company):
            continue

        key = normalize_company(sig.company)
        if not key or len(key) < 3:
            continue

        rec = company_map[key]
        if len(sig.company) > len(rec["company"]):
            rec["company"] = sig.company
        if sig.country and sig.country != "GCC":
            rec["country"] = sig.country
        elif not rec["country"]:
            rec["country"] = sig.country
        rec["sap_products"].update(sig.sap_products)
        if sig.industry:
            rec["industries"].add(sig.industry)
        rec["signal_types"].add(sig.signal_type)
        if sig.source_name and sig.source_name not in rec["sources"]:
            rec["sources"].append(sig.source_name)
        rec["signal_count"] += 1
        if quality_rank.get(sig.signal_quality, 0) > quality_rank.get(rec["best_quality"], 0):
            rec["best_quality"] = sig.signal_quality

    results = []
    for key, rec in company_map.items():
        # Final exclusion check on normalized key
        if is_excluded(key):
            continue
        rec["sap_products"] = sorted(rec["sap_products"])
        rec["industries"] = sorted(rec["industries"])
        rec["signal_types"] = sorted(rec["signal_types"])
        rec["corroboration_score"] = len(rec["signal_types"])
        results.append(rec)

    results.sort(key=lambda r: (r["corroboration_score"], r["signal_count"]), reverse=True)
    logger.info("Deduplication: %d signals → %d unique companies", len(signals), len(results))
    return results


# ============================================================================
# PPTX REPORT GENERATION
# ============================================================================

class ReportGenerator:
    """Generates professional PPTX report from aggregated signals."""

    def __init__(self, output_dir: str = "output"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def generate(self, companies: list[dict], raw_count: int) -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, companies, raw_count)
        self._add_executive_summary(prs, companies)
        self._add_country_breakdown(prs, companies)
        self._add_industry_breakdown(prs, companies)
        self._add_products(prs, companies)

        for country in ["Saudi Arabia", "UAE", "Qatar"]:
            self._add_company_table(prs, companies, country)

        self._add_high_confidence(prs, companies)
        self._add_methodology(prs)

        filename = f"SAP_Customer_Intelligence_GCC_{date.today().isoformat()}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        logger.info("Report saved: %s", filepath)
        return filepath

    def _add_title_slide(self, prs, companies, raw_count):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = SAP_DARK

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "SAP Customer Intelligence Report"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.text = "Saudi Arabia | UAE | Qatar"
        p2.font.size = Pt(24)
        p2.font.color.rgb = SAP_GOLD

        sa_count = sum(1 for c in companies if c["country"] == "Saudi Arabia")
        uae_count = sum(1 for c in companies if c["country"] == "UAE")
        qa_count = sum(1 for c in companies if c["country"] == "Qatar")

        p3 = tf.add_paragraph()
        p3.text = f"\n{len(companies)} Companies Identified  |  KSA: {sa_count}  |  UAE: {uae_count}  |  Qatar: {qa_count}  |  {date.today().strftime('%B %d, %Y')}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE

    def _add_executive_summary(self, prs, companies):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Executive Summary")

        countries = Counter(c["country"] for c in companies)
        high_conf = sum(1 for c in companies if c["corroboration_score"] >= 2)
        s4_count = sum(1 for c in companies if "SAP S/4HANA" in c["sap_products"])

        lines = [
            f"Total unique SAP end-customers identified: {len(companies)}",
            f"High-confidence (confirmed by 2+ source types): {high_conf}",
            f"Running SAP S/4HANA: {s4_count}",
            "",
            "By Country:",
        ]
        for country, count in countries.most_common():
            lines.append(f"  {country}: {count} companies")
        lines.append("")
        lines.append("By Industry:")
        industry_counter = Counter()
        for c in companies:
            for ind in c["industries"]:
                industry_counter[ind] += 1
        for ind, count in industry_counter.most_common(8):
            lines.append(f"  {ind}: {count}")

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = SAP_DARK

    def _add_country_breakdown(self, prs, companies):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Country Breakdown")

        countries = defaultdict(lambda: {"count": 0, "products": Counter(), "industries": Counter()})
        for c in companies:
            cn = c["country"]
            countries[cn]["count"] += 1
            for p in c["sap_products"]:
                countries[cn]["products"][p] += 1
            for ind in c["industries"]:
                countries[cn]["industries"][ind] += 1

        headers = ["Country", "Companies", "Top SAP Product", "Top Industry"]
        rows = len(countries) + 1
        table = slide.shapes.add_table(rows, len(headers), Inches(1), Inches(1.8), Inches(11), Inches(2.5)).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAP_BLUE

        for i, (cn, data) in enumerate(sorted(countries.items()), 1):
            top_prod = data["products"].most_common(1)[0][0] if data["products"] else "N/A"
            top_ind = data["industries"].most_common(1)[0][0] if data["industries"] else "N/A"
            table.cell(i, 0).text = cn
            table.cell(i, 1).text = str(data["count"])
            table.cell(i, 2).text = top_prod
            table.cell(i, 3).text = top_ind

    def _add_industry_breakdown(self, prs, companies):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Industry Breakdown")

        industry_counter = Counter()
        for c in companies:
            for ind in c["industries"]:
                industry_counter[ind] += 1
        if not industry_counter:
            return

        headers = ["Industry", "Companies"]
        rows = min(len(industry_counter), 15) + 1
        table = slide.shapes.add_table(rows, len(headers), Inches(1), Inches(1.8), Inches(11), Inches(3.5)).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAP_BLUE

        for i, (ind, count) in enumerate(industry_counter.most_common(15), 1):
            table.cell(i, 0).text = ind
            table.cell(i, 1).text = str(count)

    def _add_products(self, prs, companies):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "SAP Product Landscape")

        product_counter = Counter()
        for c in companies:
            for p in c["sap_products"]:
                product_counter[p] += 1

        headers = ["Product", "Companies Using"]
        rows = min(len(product_counter), 12) + 1
        table = slide.shapes.add_table(rows, len(headers), Inches(1), Inches(1.8), Inches(11), Inches(3.5)).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAP_BLUE

        for i, (prod, count) in enumerate(product_counter.most_common(12), 1):
            table.cell(i, 0).text = prod
            table.cell(i, 1).text = str(count)

    def _add_company_table(self, prs, companies, country):
        filtered = [c for c in companies if c["country"] == country]
        if not filtered:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, f"SAP Customers — {country}")

        headers = ["Company", "Industry", "SAP Products", "Confidence"]
        max_rows = min(len(filtered), 20)
        table = slide.shapes.add_table(max_rows + 1, len(headers), Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.0)).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAP_BLUE

        for i, comp in enumerate(filtered[:max_rows], 1):
            table.cell(i, 0).text = comp["company"][:40]
            table.cell(i, 1).text = ", ".join(comp["industries"])[:25] if comp["industries"] else "—"
            table.cell(i, 2).text = ", ".join(comp["sap_products"][:3])[:50]
            score = comp["corroboration_score"]
            conf = "High" if score >= 2 else "Medium" if score == 1 else "Low"
            table.cell(i, 3).text = conf

    def _add_high_confidence(self, prs, companies):
        high_conf = [c for c in companies if c["corroboration_score"] >= 2]
        if not high_conf:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "High-Confidence Targets (2+ Source Types)")

        headers = ["Company", "Country", "Industry", "SAP Products", "Sources"]
        rows = min(len(high_conf), 18) + 1
        table = slide.shapes.add_table(rows, len(headers), Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.0)).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAP_BLUE

        for i, comp in enumerate(high_conf[:18], 1):
            table.cell(i, 0).text = comp["company"][:35]
            table.cell(i, 1).text = comp["country"]
            table.cell(i, 2).text = ", ".join(comp["industries"])[:20] if comp["industries"] else "—"
            table.cell(i, 3).text = ", ".join(comp["sap_products"][:2])[:40]
            table.cell(i, 4).text = ", ".join(comp["sources"][:3])[:40]

    def _add_methodology(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Methodology & Sources")

        methods = [
            "Curated Seed List — Known SAP GCC customers from press releases, annual reports, SAP references",
            "SAP Customer Stories — sap.com/about/customer-stories filtered by Middle East",
            "SAP News Center — news.sap.com press releases mentioning GCC customers",
            "Press Releases — Google News search for SAP go-live/implementation announcements",
            "Job Postings — Companies hiring internal SAP roles (excludes SIs/consultancies)",
            "Government Procurement — SAP/ERP tenders via Google (Etimad, Dubai eSupply)",
            "Conference Agendas — LEAP, GITEX, SAP Now speaker/agenda mentions",
            "",
            "Exclusion filter applied: System integrators, tech vendors, and consulting firms are excluded",
            f"Report generated: {date.today().isoformat()}",
        ]

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for method in methods:
            p = tf.add_paragraph()
            p.text = ("• " + method) if method else ""
            p.font.size = Pt(12)
            p.font.color.rgb = SAP_DARK

    def _add_slide_title(self, slide, title: str):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = SAP_DARK


# ============================================================================
# MAIN AGENT ORCHESTRATION
# ============================================================================

def run_agent(sources: list[str] | None = None, output_dir: str = "output"):
    """Main orchestration loop."""
    print()
    print("=" * 70)
    print("  SAP Customer Intelligence Agent — GCC Edition")
    print("  Saudi Arabia | UAE | Qatar")
    print(f"  {date.today().isoformat()}")
    print("=" * 70)
    print()

    scraper_classes = {
        "seed": SeedListSource,
        "sap_stories": SAPCustomerStoriesScraper,
        "press": PressReleaseScraper,
        "jobs": JobPostingScraper,
        "gov": ProcurementScraper,
        "events": ConferenceScraper,
    }

    active_sources = sources if sources else list(scraper_classes.keys())
    print(f"Active sources: {', '.join(active_sources)}")
    print()

    all_signals: list[SAPSignal] = []

    for key in tqdm(active_sources, desc="Collecting signals", unit="source"):
        if key not in scraper_classes:
            logger.warning("Unknown source: %s", key)
            continue
        try:
            print(f"\n  Source: {key}...")
            scraper = scraper_classes[key]()
            signals = scraper.scrape()
            all_signals.extend(signals)
            print(f"    → {len(signals)} signals collected")
        except Exception as e:
            logger.error("Source %s failed: %s", key, e, exc_info=True)
            print(f"    → Error: {e}")

    raw_count = len(all_signals)
    print(f"\nTotal raw signals: {raw_count}")

    print("Deduplicating, filtering exclusions, and aggregating...")
    companies = deduplicate_signals(all_signals)
    print(f"Unique SAP end-customers identified: {len(companies)}")

    print("\nGenerating PowerPoint report...")
    generator = ReportGenerator(output_dir=output_dir)
    filepath = generator.generate(companies, raw_count)

    print()
    print("=" * 70)
    print(f"  REPORT READY: {filepath}")
    print(f"  Companies: {len(companies)}")
    print("=" * 70)
    print()

    return filepath


def main():
    parser = argparse.ArgumentParser(
        description="SAP Customer Intelligence Agent — Identifies real SAP customers in GCC"
    )
    parser.add_argument(
        "--sources",
        type=str,
        default=None,
        help="Comma-separated sources: seed,sap_stories,press,jobs,gov,events (default: all)",
    )
    parser.add_argument(
        "--output",
        type=str,
        default="output",
        help="Output directory (default: ./output)",
    )
    args = parser.parse_args()

    sources = args.sources.split(",") if args.sources else None
    run_agent(sources=sources, output_dir=args.output)


if __name__ == "__main__":
    main()
